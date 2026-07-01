#!/usr/bin/env python3
"""
build_ppt.py — 将结构化 Markdown 编译为 PPT 演示文稿。

用法：
    python build_ppt.py [markdown_file]

读取指定的 Markdown 文件，按 --- 分隔页面，识别封面页、内容页、
图片展示页，生成 output.pptx。若未指定文件，默认读取 content.md；
若 content.md 也不存在，则报错退出。
"""

import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

# ── 常量 ────────────────────────────────────────────────────

SLIDE_WIDTH  = Inches(13.333)   # 16:9 宽屏
SLIDE_HEIGHT = Inches(7.5)
BG_COVER     = "cover.png"       # 封面背景
BG_CONTENT   = "background.png"       # 内容页 / 图片页背景
INPUT_FILE   = "content.md"

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x33, 0x33, 0x33)

# 字体
FONT_LATIN = "Times New Roman"   # 英文/拉丁字体
FONT_CN_H1 = "SimSun"            # 宋体（一级标题）
FONT_CN_H2 = "SimSun"            # 宋体（二级标题）
FONT_CN_BODY = "FangSong"        # 仿宋（正文）


# ── 工具函数 ────────────────────────────────────────────────

def _set_font(p, cn_font: str, size, bold: bool = False):
    """设置段落的中文字体、英文字体、字号和粗细。"""
    p.font.name = FONT_LATIN
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = DARK

    # 通过 XML 设置东亚字体（中文字体）
    pPr = p._p.get_or_add_pPr()
    defRPr = pPr.get_or_add_defRPr()
    ea = defRPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(defRPr, qn("a:ea"))
    ea.set("typeface", cn_font)


def _clear_paragraph(paragraph):
    """清除段落中所有 run 元素。"""
    p_elem = paragraph._p
    for tag in ("a:r", "a:fld", "a:br"):
        for el in p_elem.findall(qn(tag)):
            p_elem.remove(el)


def _set_run_font(run, cn_font: str, size, bold: bool = False):
    """设置单个 run 的字体属性。"""
    run.font.name = FONT_LATIN
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = DARK
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", cn_font)


def _set_paragraph_text(paragraph, text: str, cn_font: str, size,
                        bold_default: bool = False):
    """设置段落文本，支持 **加粗** 行内语法。"""
    _clear_paragraph(paragraph)
    # 按 **...** 分割，捕获分隔符
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            inner = part[2:-2]
            if inner:
                run = paragraph.add_run()
                run.text = inner
                _set_run_font(run, cn_font, size, bold=True)
        else:
            run = paragraph.add_run()
            run.text = part
            _set_run_font(run, cn_font, size, bold=bold_default)


def _set_cell_text(cell, text: str, cn_font: str, size,
                   bold_default: bool = False):
    """设置表格单元格文本，支持 **加粗** 行内语法。"""
    tf = cell.text_frame
    p = tf.paragraphs[0]
    _set_paragraph_text(p, text, cn_font, size, bold_default)


def _set_cell_font(cell, cn_font: str, size, bold: bool = False):
    """设置表格单元格的中英文字体、字号和粗细。"""
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT_LATIN
        paragraph.font.size = size
        paragraph.font.bold = bold
        paragraph.font.color.rgb = DARK
        pPr = paragraph._p.get_or_add_pPr()
        defRPr = pPr.get_or_add_defRPr()
        ea = defRPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(defRPr, qn("a:ea"))
        ea.set("typeface", cn_font)


def _style_table_three_line(tbl, num_rows: int, num_cols: int):
    """应用三段式表格：无填充、无竖线，仅顶/中/底三条横线。"""
    BORDER_W = "19000"   # EMU ≈ 1.5pt
    BORDER_C = "333333"

    # 移除表格默认样式
    tblPr = tbl._tbl.tblPr
    for sid in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(sid)

    for ri in range(num_rows):
        for ci in range(num_cols):
            cell = tbl.cell(ri, ci)
            tcPr = cell._tc.get_or_add_tcPr()

            # 移除所有填充
            for fe in tcPr.findall(qn("a:solidFill")):
                tcPr.remove(fe)

            # 移除所有旧边框
            for tag in ("a:lnT", "a:lnB", "a:lnL", "a:lnR"):
                for el in tcPr.findall(qn(tag)):
                    tcPr.remove(el)

            # 表头行：顶线 + 底线（中间分隔线）
            if ri == 0:
                _add_tc_border(tcPr, "a:lnT", BORDER_W, BORDER_C)
                _add_tc_border(tcPr, "a:lnB", BORDER_W, BORDER_C)

            # 最后一行：底线
            if ri == num_rows - 1:
                _add_tc_border(tcPr, "a:lnB", BORDER_W, BORDER_C)


def _add_tc_border(tcPr, tag: str, width: str, color: str):
    """向单元格属性添加一条边框线。"""
    ln = etree.SubElement(tcPr, qn(tag))
    ln.set("w", width)
    sf = etree.SubElement(ln, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", color)


def _set_background(slide, image_path: str):
    """设置幻灯片背景：图片存在则铺满全幅，否则纯白。"""
    if os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        # 将背景图片移至 z-order 最底层
        sp = pic._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)
    else:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE


# ── Markdown 解析 ───────────────────────────────────────────

def parse_metadata(md_text: str) -> dict[str, str]:
    """从 Markdown 文本中解析元数据。

    支持的注释格式（须在第一个 --- 之前）：
        <!-- cover: path/to/cover.png -->
        <!-- background: path/to/bg.png -->

    返回 {"cover": "...", "background": "..."}，未指定则为空。
    """
    meta: dict[str, str] = {}
    for line in md_text.split("\n"):
        s = line.strip()
        # 遇到第一个页面分隔符即停止
        if s == "---":
            break
        m = re.match(r"<!--\s*(\w+)\s*:\s*(.+?)\s*-->", s)
        if m:
            meta[m.group(1)] = m.group(2).strip()
    return meta


def split_sections(md_text: str) -> list[str]:
    """按行首的 --- 分隔为多个页面段落。"""
    lines = md_text.split("\n")
    sections: list[str] = []
    current: list[str] = []

    for line in lines:
        if line.strip() == "---":
            if current:
                sections.append("\n".join(current).strip())
                current = []
        else:
            current.append(line)

    if current:
        sections.append("\n".join(current).strip())

    return [s for s in sections if s]


def detect_page_type(section: str) -> str:
    """检测页面类型。"""
    has_image = False
    has_structured = False
    for line in section.split("\n"):
        s = line.strip()
        if re.match(r"^#\s+", s) and not s.startswith("## "):  # 一级标题 → 封面
            return "cover"
        if re.match(r"^!\[.*\]\(.+\)", s):  # 图片引用
            has_image = True
        # 结构化内容：二级标题、列表、表格
        if (s.startswith("## ") or re.match(r"^(\d+)\.\s+", s) or
                re.match(r"^[-*]\s+", s) or (s.startswith("|") and s.endswith("|"))):
            has_structured = True

    if has_image and has_structured:
        return "text_image"
    if has_image:
        return "image"
    return "content"


def parse_cover(section: str) -> tuple[str, str]:
    """解析封面页，返回 (title, speaker)。"""
    title = ""
    speaker = ""
    for line in section.split("\n"):
        s = line.strip()
        if s.startswith("# ") and not s.startswith("## "):
            title = s[2:].strip()
        elif s.startswith("## "):
            speaker = s[3:].strip()
    return title, speaker


def parse_content(section: str) -> list[tuple]:
    """解析内容页，返回 [(type, text/headers, ...), ...]。

    type: h2 | ordered | unordered | plain | table
    表格项格式: ("table", headers_list, rows_list)
    """
    items: list[tuple] = []
    lines = section.split("\n")
    i = 0

    while i < len(lines):
        s = lines[i].strip()
        if not s:
            i += 1
            continue

        # ── 表格检测：以 | 开头和结尾的连续行 ──
        if s.startswith("|") and s.endswith("|"):
            table_lines = [s]
            j = i + 1
            while j < len(lines):
                nxt = lines[j].strip()
                if nxt.startswith("|") and nxt.endswith("|"):
                    table_lines.append(nxt)
                    j += 1
                else:
                    break

            # 至少需要 header + separator（分隔线）两行
            if len(table_lines) >= 2:
                # 解析表头（跳过首尾空元素）
                headers = [c.strip() for c in table_lines[0].split("|")[1:-1]]
                # 跳过分隔线（index 1），解析数据行
                rows = []
                for tl in table_lines[2:]:
                    row = [c.strip() for c in tl.split("|")[1:-1]]
                    if any(row):  # 跳过全空行
                        rows.append(row)
                items.append(("table", headers, rows))
                i = j
                continue

        # ── 非表格内容 ──
        # 跳过封面专属的一级标题
        if re.match(r"^#\s+", s) and not re.match(r"^##\s+", s):
            i += 1
            continue

        if s.startswith("## "):
            items.append(("h2", s[3:].strip()))
            i += 1
            continue

        m = re.match(r"^(\d+)\.\s+(.+)$", s)
        if m:
            items.append(("ordered", m.group(2).strip(), int(m.group(1))))
            i += 1
            continue

        m = re.match(r"^[-*]\s+(.+)$", s)
        if m:
            items.append(("unordered", m.group(1).strip()))
            i += 1
            continue

        items.append(("plain", s))
        i += 1

    return items


def parse_image_page(section: str) -> tuple[str, list[tuple[str, str]]]:
    """解析图片展示页，返回 (caption, [(alt_text, path), ...])。"""
    caption = ""
    images: list[tuple[str, str]] = []
    for line in section.split("\n"):
        s = line.strip()
        if not s:
            continue
        m = re.match(r"^!\[(.*)\]\((.+)\)$", s)
        if m:
            alt = m.group(1).strip()
            path = m.group(2).strip()
            images.append((alt, path))
        elif not s.startswith("#"):
            caption = s
    return caption, images


def parse_text_image(section: str) -> tuple[list[tuple], str]:
    """解析图配文页，返回 (items, image_path)。
    文字部分解析逻辑同 content 页，图片引用行单独提取。
    """
    items: list[tuple] = []
    image_path = ""
    lines = section.split("\n")
    i = 0

    while i < len(lines):
        s = lines[i].strip()
        if not s:
            i += 1
            continue

        # 图片引用行 → 提取路径，不加入 items
        m = re.match(r"^!\[.*\]\((.+)\)$", s)
        if m:
            image_path = m.group(1).strip()
            i += 1
            continue

        # ── 表格检测：以 | 开头和结尾的连续行 ──
        if s.startswith("|") and s.endswith("|"):
            table_lines = [s]
            j = i + 1
            while j < len(lines):
                nxt = lines[j].strip()
                if nxt.startswith("|") and nxt.endswith("|"):
                    table_lines.append(nxt)
                    j += 1
                else:
                    break

            if len(table_lines) >= 2:
                headers = [c.strip() for c in table_lines[0].split("|")[1:-1]]
                rows = []
                for tl in table_lines[2:]:
                    row = [c.strip() for c in tl.split("|")[1:-1]]
                    if any(row):
                        rows.append(row)
                items.append(("table", headers, rows))
                i = j
                continue

        # 跳过封面专属的一级标题
        if re.match(r"^#\s+", s) and not re.match(r"^##\s+", s):
            i += 1
            continue

        if s.startswith("## "):
            items.append(("h2", s[3:].strip()))
            i += 1
            continue

        m2 = re.match(r"^(\d+)\.\s+(.+)$", s)
        if m2:
            items.append(("ordered", m2.group(2).strip(), int(m2.group(1))))
            i += 1
            continue

        m2 = re.match(r"^[-*]\s+(.+)$", s)
        if m2:
            items.append(("unordered", m2.group(1).strip()))
            i += 1
            continue

        items.append(("plain", s))
        i += 1

    return items, image_path


def _fit_image_size(img_path: str, max_w: float, max_h: float) -> tuple[float, float]:
    """根据图片原始尺寸，在 max_w × max_h 范围内等比例缩放，返回 (w, h)。"""
    with Image.open(img_path) as im:
        orig_w, orig_h = im.size
    scale = min(max_w / orig_w, max_h / orig_h, 1.0)  # 不放大小图
    return orig_w * scale, orig_h * scale


# ── 幻灯片创建 ──────────────────────────────────────────────

def _new_blank_slide(prs: Presentation):
    """创建一个空白版式的幻灯片。"""
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


def create_cover_slide(prs: Presentation, title: str, speaker: str):
    """创建封面幻灯片。"""
    slide = _new_blank_slide(prs)
    _set_background(slide, BG_COVER)

    if title:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_paragraph_text(p, title, FONT_CN_H1, Pt(66), bold_default=True)

    if speaker:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.333), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_paragraph_text(p, speaker, FONT_CN_H2, Pt(44))


def create_content_slide(prs: Presentation, items: list[tuple]):
    """创建内容幻灯片，支持文本和表格。"""
    slide = _new_blank_slide(prs)
    _set_background(slide, BG_CONTENT)

    if not items:
        return

    # 分离文本项和表格项
    text_items = [it for it in items if it[0] != "table"]
    table_items = [it for it in items if it[0] == "table"]

    # ── 特殊处理：仅有一个二级标题且无表格 → 居中显示 ──
    if (len(text_items) == 1 and text_items[0][0] == "h2"
            and not table_items):
        # 文本框居中放置，水平和垂直均居中
        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.75), Inches(10.333), Inches(2.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_paragraph_text(p, text_items[0][1], FONT_CN_H2, Pt(44), bold_default=True)
        return

    # ── 渲染文本 ──
    if text_items:
        # 有表格时文本框高度缩小，为表格留空间
        text_height = Inches(2.5) if table_items else Inches(5.8)
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10.333), text_height)
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for item in text_items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            t = item[0]

            if t == "h2":
                _set_paragraph_text(p, item[1], FONT_CN_H2, Pt(44), bold_default=True)
                p.space_after = Pt(24)
            elif t == "ordered":
                _set_paragraph_text(p, f"{item[2]}. {item[1]}", FONT_CN_BODY, Pt(32))
                p.space_after = Pt(12)
            elif t == "unordered":
                _set_paragraph_text(p, f"• {item[1]}", FONT_CN_BODY, Pt(32))
                p.space_after = Pt(12)
            else:  # plain
                _set_paragraph_text(p, item[1], FONT_CN_BODY, Pt(32))
                p.space_after = Pt(12)

    # ── 渲染表格 ──
    if table_items:
        table_top = Inches(3.5) if text_items else Inches(0.8)
        for _, headers, rows in table_items:
            num_cols = len(headers)
            num_rows = len(rows) + 1  # +1 for header

            col_width = Inches(10.333 / max(num_cols, 1))
            row_height = Inches(0.5)

            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(1.5), table_top,
                col_width * num_cols, row_height * num_rows
            )
            tbl = table_shape.table

            # 设置列宽
            for ci in range(num_cols):
                tbl.columns[ci].width = col_width

            # 表头行
            for ci, header_text in enumerate(headers):
                cell = tbl.cell(0, ci)
                _set_cell_text(cell, header_text, FONT_CN_H2, Pt(32), bold_default=True)

            # 数据行
            for ri, row in enumerate(rows):
                for ci, cell_text in enumerate(row):
                    if ci < num_cols:
                        cell = tbl.cell(ri + 1, ci)
                        _set_cell_text(cell, cell_text, FONT_CN_BODY, Pt(32))

            # 应用三段式样式（无填充、仅三条横线）
            _style_table_three_line(tbl, num_rows, num_cols)

            # 后续表格依次向下排列
            table_top += row_height * num_rows + Inches(0.3)


def create_image_slide(prs: Presentation, caption: str,
                        images: list[tuple[str, str]]):
    """创建图片展示幻灯片，支持单图/多图自动布局。

    布局规则：
    - 1 张：居中大幅展示
    - 2 张：左右均分
    - 3 张：上排 2 张 + 下排 1 张居中
    - 4 张：2×2 田字格
    - 5～6 张：3 列网格，最多 2 行
    - 超过 6 张：仅取前 6 张并警告
    """
    slide = _new_blank_slide(prs)
    _set_background(slide, BG_CONTENT)

    # 过滤存在的图片
    valid = [(alt, p) for alt, p in images if os.path.exists(p)]
    for alt, p in images:
        if not os.path.exists(p):
            print(f"  ⚠ 图片不存在: {p}")

    MAX_IMAGES = 6
    if len(valid) > MAX_IMAGES:
        print(f"  ⚠ 图片过多 ({len(valid)} 张)，仅显示前 {MAX_IMAGES} 张")
        valid = valid[:MAX_IMAGES]

    if not valid:
        if caption:
            tb = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _set_paragraph_text(p, caption, FONT_CN_H2, Pt(44), bold_default=True)
        return

    n = len(valid)

    # ── 标题区域 ──
    CAP_TOP = 0.3
    CAP_H = 0.9
    if caption:
        tb = slide.shapes.add_textbox(
            Inches(1), Inches(CAP_TOP), Inches(11.333), Inches(CAP_H))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_paragraph_text(p, caption, FONT_CN_H2, Pt(44), bold_default=True)

    # ── 图片可用区域 ──
    MARGIN = 0.8
    GAP = 0.15
    area_left = MARGIN
    area_top = CAP_TOP + CAP_H + 0.15 if caption else 0.5
    area_w = 13.333 - 2 * MARGIN
    area_h = 7.5 - area_top - 0.3

    if n == 1:
        _place_single_image(slide, valid[0], area_left, area_top, area_w, area_h)
    elif n == 2:
        _place_image_grid(slide, valid, area_left, area_top, area_w, area_h,
                          cols=2, rows=1, gap=GAP)
    elif n == 3:
        row_h = (area_h - GAP) / 2
        col_w2 = (area_w - GAP) / 2
        # 上排 2 张
        for ci in range(2):
            cx = area_left + ci * (col_w2 + GAP)
            _place_cell(slide, valid[ci], cx, area_top, col_w2, row_h)
        # 下排 1 张居中
        col_w1 = area_w * 0.5
        cx = area_left + (area_w - col_w1) / 2
        _place_cell(slide, valid[2], cx, area_top + row_h + GAP, col_w1, row_h)
    elif n == 4:
        _place_image_grid(slide, valid, area_left, area_top, area_w, area_h,
                          cols=2, rows=2, gap=GAP)
    else:  # 5～6
        _place_image_grid(slide, valid, area_left, area_top, area_w, area_h,
                          cols=3, rows=2, gap=GAP)


def _place_cell(slide, image: tuple[str, str],
                left: float, top: float, cell_w: float, cell_h: float):
    """在指定单元格内居中放置一张图片，alt 文本显示在图片下方。"""
    alt_text, path = image
    SUB_H = 0.28 if alt_text else 0.0
    img_h = cell_h - SUB_H - 0.04

    pic_w, pic_h = _fit_image_size(path, cell_w - 0.08, img_h - 0.08)
    pic_left = left + (cell_w - pic_w) / 2
    pic_top = top + (img_h - pic_h) / 2
    slide.shapes.add_picture(path, Inches(pic_left), Inches(pic_top),
                             Inches(pic_w), Inches(pic_h))

    if alt_text:
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top + cell_h - SUB_H),
            Inches(cell_w), Inches(SUB_H))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_paragraph_text(p, alt_text, FONT_CN_BODY, Pt(14))


def _place_single_image(slide, image: tuple[str, str],
                        left: float, top: float, area_w: float, area_h: float):
    """单图居中，保留原有 8.0×4.8 的舒适展示区域。"""
    alt_text, path = image
    max_w, max_h = 8.0, 4.8
    pic_w, pic_h = _fit_image_size(path, max_w, max_h)
    pic_left = left + (area_w - pic_w) / 2
    pic_top = top + (area_h - pic_h) / 2
    slide.shapes.add_picture(path, Inches(pic_left), Inches(pic_top),
                             Inches(pic_w), Inches(pic_h))

    if alt_text:
        sub_top = pic_top + pic_h + 0.1
        tb = slide.shapes.add_textbox(
            Inches(pic_left), Inches(sub_top),
            Inches(pic_w), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_paragraph_text(p, alt_text, FONT_CN_BODY, Pt(16))


def _place_image_grid(slide, images: list[tuple[str, str]],
                      left: float, top: float, area_w: float, area_h: float,
                      cols: int, rows: int, gap: float):
    """将图片按 cols×rows 网格排列。"""
    col_w = (area_w - (cols - 1) * gap) / cols
    row_h = (area_h - (rows - 1) * gap) / rows

    for idx, img in enumerate(images):
        if idx >= cols * rows:
            break
        r = idx // cols
        c = idx % cols
        cx = left + c * (col_w + gap)
        cy = top + r * (row_h + gap)
        _place_cell(slide, img, cx, cy, col_w, row_h)


def create_text_image_slide(prs: Presentation, items: list[tuple], img_path: str):
    """创建图配文幻灯片：左侧 2/3 文字，右侧 1/3 图片。
    文字部分支持的数据类型和样式与原"内容页"一致。
    """
    slide = _new_blank_slide(prs)
    _set_background(slide, BG_CONTENT)

    # ── 左侧 2/3：文字区域 ──
    text_left = Inches(0.8)
    text_width = Inches(8.0)

    text_items = [it for it in items if it[0] != "table"]
    table_items = [it for it in items if it[0] == "table"]

    if text_items:
        text_height = Inches(2.5) if table_items else Inches(5.8)
        tb = slide.shapes.add_textbox(text_left, Inches(0.8), text_width, text_height)
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for item in text_items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            t = item[0]

            if t == "h2":
                _set_paragraph_text(p, item[1], FONT_CN_H2, Pt(44), bold_default=True)
                p.space_after = Pt(24)
            elif t == "ordered":
                _set_paragraph_text(p, f"{item[2]}. {item[1]}", FONT_CN_BODY, Pt(32))
                p.space_after = Pt(12)
            elif t == "unordered":
                _set_paragraph_text(p, f"• {item[1]}", FONT_CN_BODY, Pt(32))
                p.space_after = Pt(12)
            else:  # plain
                _set_paragraph_text(p, item[1], FONT_CN_BODY, Pt(32))
                p.space_after = Pt(12)

    if table_items:
        table_top = Inches(3.5) if text_items else Inches(0.8)
        for _, headers, rows in table_items:
            num_cols = len(headers)
            num_rows = len(rows) + 1

            col_width = Inches(text_width.inches / max(num_cols, 1))
            row_height = Inches(0.5)

            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                text_left, table_top,
                col_width * num_cols, row_height * num_rows
            )
            tbl = table_shape.table

            for ci in range(num_cols):
                tbl.columns[ci].width = col_width

            for ci, header_text in enumerate(headers):
                cell = tbl.cell(0, ci)
                _set_cell_text(cell, header_text, FONT_CN_H2, Pt(32), bold_default=True)

            for ri, row in enumerate(rows):
                for ci, cell_text in enumerate(row):
                    if ci < num_cols:
                        cell = tbl.cell(ri + 1, ci)
                        _set_cell_text(cell, cell_text, FONT_CN_BODY, Pt(32))

            _style_table_three_line(tbl, num_rows, num_cols)
            table_top += row_height * num_rows + Inches(0.3)

    # ── 右侧 1/3：图片区域（等比例缩放，维持原图比例）──
    if img_path and os.path.exists(img_path):
        img_left = Inches(9.2)
        max_w, max_h = 3.6, 5.5  # 英寸
        pic_w, pic_h = _fit_image_size(img_path, max_w, max_h)
        # 在右侧区域内居中
        pic_left = 9.2 + (max_w - pic_w) / 2
        pic_top = 1.5 + (max_h - pic_h) / 2
        slide.shapes.add_picture(img_path, Inches(pic_left), Inches(pic_top),
                                 Inches(pic_w), Inches(pic_h))
    elif img_path:
        print(f"  ⚠ 图片不存在: {img_path}")


# ── 主流程 ──────────────────────────────────────────────────

def create_end_slide(prs: Presentation):
    """创建结束页：「谢谢」。"""
    slide = _new_blank_slide(prs)
    _set_background(slide, BG_COVER)

    tb = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.333), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "谢谢"
    p.alignment = PP_ALIGN.CENTER
    _set_font(p, FONT_CN_H1, Pt(66), bold=True)


def build(md_path: str = INPUT_FILE):
    """读取 Markdown 文件，在同目录生成同名 .pptx。"""
    global BG_COVER, BG_CONTENT

    if not os.path.exists(md_path):
        print(f"错误：找不到输入文件 {md_path}")
        sys.exit(1)

    md_abs = os.path.abspath(md_path)
    md_dir = os.path.dirname(md_abs)
    output_path = os.path.splitext(md_abs)[0] + ".pptx"

    with open(md_path, "r", encoding="utf-8") as f:
        md_text = f.read()

    # 解析元数据，覆盖默认背景图
    meta = parse_metadata(md_text)
    def _resolve_bg(key: str, default: str) -> str:
        """解析背景图路径：元数据优先，相对路径基于 md 文件目录。"""
        path = meta.get(key) or default
        if not os.path.isabs(path):
            path = os.path.normpath(os.path.join(md_dir, path))
        return path

    BG_COVER = _resolve_bg("cover", "cover.png")
    BG_CONTENT = _resolve_bg("background", "background.png")

    sections = split_sections(md_text)
    if not sections:
        print("错误：Markdown 中没有有效内容")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, sec in enumerate(sections):
        pt = detect_page_type(sec)
        print(f"  第 {i + 1} 页 → {pt}")

        if pt == "cover":
            title, speaker = parse_cover(sec)
            create_cover_slide(prs, title, speaker)
        elif pt == "image":
            caption, images = parse_image_page(sec)
            create_image_slide(prs, caption, images)
        elif pt == "text_image":
            items, img_path = parse_text_image(sec)
            create_text_image_slide(prs, items, img_path)
        else:
            items = parse_content(sec)
            create_content_slide(prs, items)

    # 自动添加结束页
    create_end_slide(prs)
    print(f"  第 {len(sections) + 1} 页 → end")

    prs.save(output_path)
    print(f"\n已生成 {output_path}（共 {len(sections)} 页）")


if __name__ == "__main__":
    if len(sys.argv) > 1:
        md_path = sys.argv[1]
    elif os.path.exists(INPUT_FILE):
        md_path = INPUT_FILE
    else:
        print(f"错误：未指定输入文件，且默认文件 {INPUT_FILE} 不存在")
        print(f"用法：python build_ppt.py [markdown_file]")
        sys.exit(1)
    build(md_path)
